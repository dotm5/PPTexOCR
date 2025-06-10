import sys
import os
import io
import ctypes
from PySide6.QtWidgets import (
    QApplication, QWidget, QVBoxLayout, QHBoxLayout,
    QPushButton, QListWidget, QListWidgetItem, QFileDialog,
    QMessageBox, QLabel, QTextEdit, QProgressBar, QSpacerItem, QSizePolicy
)
from PySide6.QtCore import Qt, QThread, Signal, QPoint, QRectF
from PySide6.QtGui import (
    QFont, QColor, QPainter, QBrush, QPen, QIcon, QAction, QCursor,
    QPainterPath, QRegion
)
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import pytesseract
import torch
from pix2tex.cli import LatexOCR


# 初始化pix2tex LatexOCR模型
pix2tex_model = LatexOCR()
device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')


def enable_blur_behind_window(win):
    """Windows平台启用毛玻璃效果"""
    if sys.platform != "win32":
        return

    hwnd = int(win.winId())

    class DWM_BLURBEHIND(ctypes.Structure):
        _fields_ = [
            ("dwFlags", ctypes.c_uint),
            ("fEnable", ctypes.c_int),
            ("hRgnBlur", ctypes.c_void_p),
            ("fTransitionOnMaximized", ctypes.c_int),
        ]

    DWM_BB_ENABLE = 0x1
    dwmapi = ctypes.windll.dwmapi
    blur_behind = DWM_BLURBEHIND()
    blur_behind.dwFlags = DWM_BB_ENABLE
    blur_behind.fEnable = 1
    blur_behind.hRgnBlur = None
    blur_behind.fTransitionOnMaximized = 0

    dwmapi.DwmEnableBlurBehindWindow(hwnd, ctypes.byref(blur_behind))


def ocr_image_multilang_with_latex(img: Image.Image) -> str:
    """对图片进行OCR和LaTeX公式识别"""
    try:
        text_ocr = pytesseract.image_to_string(img, lang='eng+chi_sim')
    except Exception:
        text_ocr = ""

    try:
        latex_text = pix2tex_model(img)
    except Exception:
        latex_text = ""

    combined = ""
    if text_ocr.strip():
        combined += "[OCR Text]\n" + text_ocr.strip() + "\n"
    if latex_text.strip():
        combined += "[LaTeX OCR]\n" + latex_text.strip() + "\n"

    return combined.strip()


def extract_text_from_pptx(path: str) -> str:
    """提取pptx中的文本及对图片进行OCR识别"""
    prs = Presentation(path)
    all_text = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_texts = []

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or shape.has_text_frame:
                text_frame = shape.text_frame
                if text_frame:
                    texts = [p.text for p in text_frame.paragraphs if p.text.strip()]
                    slide_texts.extend(texts)

            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                try:
                    img = Image.open(io.BytesIO(image_bytes))
                    ocr_result = ocr_image_multilang_with_latex(img)
                    if ocr_result:
                        slide_texts.append(ocr_result)
                except Exception as e:
                    slide_texts.append(f"[Error OCR Image: {e}]")

        if slide_texts:
            all_text.append(f"--- Slide {slide_idx + 1} ---")
            all_text.extend(slide_texts)

    return "\n".join(all_text)


class WorkerThread(QThread):
    progress = Signal(str, str)  # filepath, status message
    finished = Signal(str, str)  # filepath, recognized text
    error = Signal(str)          # error message

    def __init__(self, filepath: str):
        super().__init__()
        self.filepath = filepath

    def run(self):
        try:
            self.progress.emit(self.filepath, "识别中...")
            text = extract_text_from_pptx(self.filepath)
            self.finished.emit(self.filepath, text)
        except Exception as e:
            self.error.emit(f"文件 {os.path.basename(self.filepath)} 识别失败: {e}")
            self.finished.emit(self.filepath, "")


class PPTOCRApp(QWidget):
    def __init__(self):
        super().__init__()

        # 无边框透明窗口，圆角化
        self.setWindowFlag(Qt.FramelessWindowHint)
        self.setAttribute(Qt.WA_TranslucentBackground)
        self.resize(900, 680)
        self.dragPos = QPoint()

        self.file_list = []  # [filepath, status, recognized_text]
        self.threads = []

        self.setAcceptDrops(True)  # 支持拖放

        self.init_ui()
        self.apply_glass_style()
        self.apply_palette_colors()
        enable_blur_behind_window(self)

        # 设置窗口圆角遮罩
        self.update_mask()

    def init_ui(self):
        main_layout = QVBoxLayout(self)
        main_layout.setContentsMargins(12, 12, 12, 12)
        main_layout.setSpacing(8)

        # 自定义标题栏
        self.title_bar = QWidget()
        self.title_bar.setFixedHeight(38)
        self.title_bar.setObjectName("titleBar")
        title_layout = QHBoxLayout(self.title_bar)
        title_layout.setContentsMargins(10, 0, 10, 0)
        title_layout.setSpacing(8)

        self.title_label = QLabel("PPT 批量文字+公式 OCR 识别（pix2tex + PySide6）")
        self.title_label.setFont(QFont("Segoe UI", 12, QFont.Weight.Bold))
        self.title_label.setStyleSheet("color: white;")
        title_layout.addWidget(self.title_label)

        title_layout.addSpacerItem(QSpacerItem(40, 20, QSizePolicy.Expanding, QSizePolicy.Minimum))

        # 关闭按钮
        self.btn_close = QPushButton('×')
        self.btn_close.setFixedSize(36, 28)
        self.btn_close.setFont(QFont("Segoe UI", 14))
        self.btn_close.setCursor(Qt.PointingHandCursor)
        self.btn_close.setToolTip("关闭")
        self.btn_close.setObjectName("btnClose")
        self.btn_close.clicked.connect(self.close)
        title_layout.addWidget(self.btn_close)

        # 最小化按钮
        self.btn_minimize = QPushButton('−')
        self.btn_minimize.setFixedSize(36, 28)
        self.btn_minimize.setFont(QFont("Segoe UI", 14))
        self.btn_minimize.setCursor(Qt.PointingHandCursor)
        self.btn_minimize.setToolTip("最小化")
        self.btn_minimize.setObjectName("btnMinimize")
        self.btn_minimize.clicked.connect(self.showMinimized)
        title_layout.addWidget(self.btn_minimize)

        # 按钮区域
        btn_layout = QHBoxLayout()
        btn_layout.setSpacing(12)

        self.btn_add = QPushButton("添加PPT文件")
        self.btn_remove = QPushButton("移除选中文件")
        self.btn_start = QPushButton("开始识别")
        self.btn_export = QPushButton("导出选中文本")

        for btn in (self.btn_add, self.btn_remove, self.btn_start, self.btn_export):
            btn.setMinimumHeight(36)
            btn.setCursor(Qt.PointingHandCursor)
            btn.setFont(QFont("Segoe UI", 11))

        btn_layout.addWidget(self.btn_add)
        btn_layout.addWidget(self.btn_remove)
        btn_layout.addWidget(self.btn_start)
        btn_layout.addWidget(self.btn_export)

        # 文件列表
        self.list_widget = QListWidget()
        self.list_widget.setSelectionMode(QListWidget.SelectionMode.ExtendedSelection)
        self.list_widget.setFont(QFont("Segoe UI", 10))

        # 状态区
        self.status_label = QLabel("准备就绪")
        self.status_label.setFont(QFont("Segoe UI", 10))

        self.progress_bar = QProgressBar()
        self.progress_bar.setMinimumHeight(20)
        self.progress_bar.setFormat("0 / 0")
        self.progress_bar.setTextVisible(True)

        # 日志输出
        self.log_text = QTextEdit()
        self.log_text.setReadOnly(True)
        self.log_text.setFont(QFont("Segoe UI", 10))

        main_layout.addWidget(self.title_bar, stretch=0)
        main_layout.addLayout(btn_layout)
        main_layout.addWidget(self.list_widget, stretch=1)
        main_layout.addWidget(self.status_label)
        main_layout.addWidget(self.progress_bar)
        main_layout.addWidget(QLabel("识别日志："))
        main_layout.addWidget(self.log_text, stretch=1)

        # 信号槽绑定
        self.btn_add.clicked.connect(self.on_add_files_button_clicked)
        self.btn_remove.clicked.connect(self.remove_selected_files)
        self.btn_start.clicked.connect(self.start_recognition)
        self.btn_export.clicked.connect(self.export_selected_texts)

    def apply_glass_style(self):
        """
        给窗口和所有子控件应用半透明玻璃效果(透明背景+半透明配色)
        通过设置WA_TranslucentBackground属性和统一样式表实现
        """
        # 统一半透明背景色和文字色，递归设置子控件透明背景
        def set_glass_style_recursive(widget):
            widget.setAttribute(Qt.WA_TranslucentBackground)
            widget.setStyleSheet("""
                background-color: rgba(30, 30, 30, 180); /* 半透明深灰背景 */
                color: #E0E0E0; /* 文字浅灰，保证对比 */
                selection-background-color: rgba(100, 100, 255, 120); /* 选中半透明蓝 */
                border-radius: 7px;
            """)
            for child in widget.findChildren(QWidget):
                set_glass_style_recursive(child)

        set_glass_style_recursive(self)

        # 特殊控件单独调整样式，覆盖统一透明背景
        self.list_widget.setStyleSheet("""
            QListWidget {
                background-color: rgba(40, 40, 40, 230);  /* 稍微更透明 */
                border: 1px solid rgba(120, 120, 120, 90);
                border-radius: 8px;
                color: #E0E0E0;
                selection-background-color: rgba(50, 100, 255, 160);
                padding: 5px;
            }
            QListWidget::item:selected {
                color: white;
                background-color: rgba(70, 130, 255, 200);
                border-radius: 4px;
            }
        """)

        self.log_text.setStyleSheet("""
            QTextEdit {
                background-color: rgba(30, 30, 30, 200);
                border: 1px solid rgba(90, 90, 90, 150);
                border-radius: 8px;
                color: #E0E0E0;
            }
        """)

        self.status_label.setStyleSheet("color: #CCCCCC; padding: 4px;")

        self.progress_bar.setStyleSheet("""
            QProgressBar {
                border: 1px solid rgba(100, 100, 100, 90);
                border-radius: 10px;
                background-color: rgba(50, 50, 50, 180);
                text-align: center;
                color: #CCCCCC;
            }
            QProgressBar::chunk {
                background-color: rgba(80, 150, 255, 200);
                border-radius: 10px;
            }
        """)

        # 按钮半透明背景及hover效果，颜色配置位置
        btn_style = """
            QPushButton {
                background-color: rgba(70, 130, 255, 180); /* 半透明蓝 */
                color: white;
                border: none;
                border-radius: 6px;
                padding: 8px 14px;
                font-weight: 600;
                transition: background-color 0.2s ease;
            }
            QPushButton:hover {
                background-color: rgba(70, 130, 255, 230); /* hover更亮 */
            }
            QPushButton:pressed {
                background-color: rgba(50, 110, 220, 230); /* 按下更暗 */
            }
        """
        for btn in (self.btn_add, self.btn_remove, self.btn_start, self.btn_export):
            btn.setStyleSheet(btn_style)

        # 标题栏按钮
        self.title_bar.setStyleSheet("""
            QWidget#titleBar {
                background-color: rgba(20, 20, 20, 200);
                border-top-left-radius: 20px;
                border-top-right-radius: 20px;
            }
        """)
        self.btn_close.setStyleSheet("""
            QPushButton#btnClose {
                background-color: rgba(220, 50, 50, 200);
                border: none;
                border-radius: 6px;
                color: white;
            }
            QPushButton#btnClose:hover {
                background-color: rgba(255, 80, 80, 255);
            }
            QPushButton#btnClose:pressed {
                background-color: rgba(180, 40, 40, 255);
            }
        """)
        self.btn_minimize.setStyleSheet("""
            QPushButton#btnMinimize {
                background-color: rgba(100, 100, 100, 150);
                border: none;
                border-radius: 6px;
                color: white;
            }
            QPushButton#btnMinimize:hover {
                background-color: rgba(140, 140, 140, 200);
            }
            QPushButton#btnMinimize:pressed {
                background-color: rgba(80, 80, 80, 200);
            }
        """)

    def apply_palette_colors(self):
        """
        调整应用整体调色板，主要字体颜色和控件前景色
        """
        palette = self.palette()
        palette.setColor(self.foregroundRole(), QColor("#E0E0E0"))  # 主文字色
        palette.setColor(self.backgroundRole(), QColor(30, 30, 30, 180))  # 背景色（半透明）
        self.setPalette(palette)

    # 支持拖放
    def dragEnterEvent(self, event):
        if event.mimeData().hasUrls():
            for url in event.mimeData().urls():
                if url.toLocalFile().lower().endswith(".pptx"):
                    event.acceptProposedAction()
                    return
        event.ignore()

    def dropEvent(self, event):
        urls = event.mimeData().urls()
        paths = [url.toLocalFile() for url in urls if url.toLocalFile().lower().endswith(".pptx")]
        self.add_files(paths)

    # 文件添加入口（无参数调用弹窗）
    def on_add_files_button_clicked(self):
        files, _ = QFileDialog.getOpenFileNames(self, "选择PPTX文件", "", "PPTX文件 (*.pptx)")
        if not files:
            return
        self.add_files(files)

    # 添加文件到列表
    def add_files(self, files):
        existing_files = {item[0] for item in self.file_list}
        count_added = 0
        for f in files:
            if f not in existing_files:
                self.file_list.append([f, "待处理", ""])
                item = QListWidgetItem(os.path.basename(f) + "  [待处理]")
                self.list_widget.addItem(item)
                count_added += 1
        if count_added:
            self.update_progress()

    def remove_selected_files(self):
        selected = self.list_widget.selectedIndexes()
        if not selected:
            return
        rows = sorted({idx.row() for idx in selected}, reverse=True)
        for row in rows:
            self.list_widget.takeItem(row)
            del self.file_list[row]
        self.update_progress()

    def start_recognition(self):
        if not self.file_list:
            QMessageBox.warning(self, "提示", "请先添加PPTX文件")
            return

        # 禁用按钮防止重复操作
        self.btn_add.setEnabled(False)
        self.btn_remove.setEnabled(False)
        self.btn_start.setEnabled(False)
        self.btn_export.setEnabled(False)

        self.log_text.clear()
        self.status_label.setText("开始识别任务...")
        self.progress_bar.setMaximum(len(self.file_list))
        self.progress_bar.setValue(0)
        self.progress_bar.setFormat(f"0 / {len(self.file_list)}")

        self.threads = []

        for idx, (filepath, status, _) in enumerate(self.file_list):
            if status == "待处理":
                worker = WorkerThread(filepath)
                worker.progress.connect(self.on_worker_progress)
                worker.finished.connect(self.on_worker_finished)
                worker.error.connect(self.on_worker_error)
                self.threads.append(worker)
                self.file_list[idx][1] = "排队中"
                self.update_list_item(idx)
                worker.start()

    def on_worker_progress(self, filepath, msg):
        idx = self.index_of_filepath(filepath)
        if idx != -1:
            self.file_list[idx][1] = msg
            self.update_list_item(idx)
            self.status_label.setText(f"{os.path.basename(filepath)}: {msg}")

    def on_worker_finished(self, filepath, text):
        idx = self.index_of_filepath(filepath)
        if idx != -1:
            if text:
                self.file_list[idx][1] = "完成"
                self.file_list[idx][2] = text
            else:
                self.file_list[idx][1] = "失败"
            self.update_list_item(idx)

        done_count = sum(1 for f in self.file_list if f[1] in ("完成", "失败"))
        self.progress_bar.setValue(done_count)
        self.progress_bar.setFormat(f"{done_count} / {len(self.file_list)}")

        if done_count == len(self.file_list):
            self.status_label.setText("所有任务完成")
            self.btn_add.setEnabled(True)
            self.btn_remove.setEnabled(True)
            self.btn_start.setEnabled(True)
            self.btn_export.setEnabled(True)

    def on_worker_error(self, msg):
        self.log_text.append(f"<span style='color:#FF6666;'>{msg}</span>")

    def update_list_item(self, idx):
        item = self.list_widget.item(idx)
        filepath, status, _ = self.file_list[idx]
        item.setText(f"{os.path.basename(filepath)}  [{status}]")

    def index_of_filepath(self, filepath):
        for i, (fp, _, _) in enumerate(self.file_list):
            if fp == filepath:
                return i
        return -1

    def export_selected_texts(self):
        selected = self.list_widget.selectedIndexes()
        if not selected:
            QMessageBox.warning(self, "提示", "请先选中至少一个文件导出")
            return

        for idx in {idx.row() for idx in selected}:
            filepath, status, text = self.file_list[idx]
            if status != "完成":
                QMessageBox.warning(self, "提示", f"文件 {os.path.basename(filepath)} 未完成识别，不能导出")
                continue
            if not text.strip():
                QMessageBox.warning(self, "提示", f"文件 {os.path.basename(filepath)} 无识别文本")
                continue

            default_name = os.path.splitext(os.path.basename(filepath))[0] + ".txt"
            save_path, _ = QFileDialog.getSaveFileName(self, "保存文本文件", default_name, "文本文件 (*.txt)")
            if save_path:
                try:
                    with open(save_path, "w", encoding="utf-8") as f:
                        f.write(text)
                    self.log_text.append(f"导出成功: {save_path}")
                except Exception as e:
                    QMessageBox.warning(self, "错误", f"保存文件失败: {e}")

    def update_progress(self):
        self.progress_bar.setMaximum(len(self.file_list))
        done_count = sum(1 for f in self.file_list if f[1] in ("完成", "失败"))
        self.progress_bar.setValue(done_count)
        self.progress_bar.setFormat(f"{done_count} / {len(self.file_list)}")

    # 重绘圆角背景
    def paintEvent(self, event):
        radius = 20  # 圆角半径
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)
        rect = self.rect()
        color = QColor(30, 30, 30, 220)  # 半透明深灰背景
        painter.setBrush(QBrush(color))
        painter.setPen(Qt.NoPen)
        painter.drawRoundedRect(QRectF(rect), radius, radius)

    # 实现拖动窗口
    def mousePressEvent(self, event):
        if event.button() == Qt.LeftButton and self.is_in_title_bar(event.pos()):
            self.dragPos = event.globalPosition().toPoint()
            event.accept()
        else:
            super().mousePressEvent(event)

    def mouseMoveEvent(self, event):
        if event.buttons() & Qt.LeftButton:
            if hasattr(self, 'dragPos'):
                self.move(self.pos() + event.globalPosition().toPoint() - self.dragPos)
                self.dragPos = event.globalPosition().toPoint()
                event.accept()
        else:
            super().mouseMoveEvent(event)

    def is_in_title_bar(self, pos):
        return self.title_bar.geometry().contains(pos)

    # 关键: 设置圆角遮罩，裁剪窗口形状避免矩形半透明区域
    def update_mask(self):
        radius = 20
        path = QPainterPath()
        path.addRoundedRect(self.rect(), radius, radius)
        region = QRegion(path.toFillPolygon().toPolygon())
        self.setMask(region)

    def resizeEvent(self, event):
        super().resizeEvent(event)
        self.update_mask()


if __name__ == "__main__":
    # 启用高DPI支持
    QApplication.setAttribute(Qt.AA_EnableHighDpiScaling)

    app = QApplication(sys.argv)
    app.setFont(QFont("Segoe UI", 10))

    window = PPTOCRApp()
    window.show()

    sys.exit(app.exec())
